"""Genera presentación .pptx documentando el QC técnico M4 de ecoRTA.

Uso:
    python scripts/generate_qc_slides.py

Salida: output/ecoRTA_QC_tecnico_M4.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Paleta de colores (inspirada en la UI arcade del overlay)
# ---------------------------------------------------------------------------

C_BG_DARK  = RGBColor(0x07, 0x05, 0x0F)   # fondo negro profundo
C_ACCENT   = RGBColor(0xA7, 0x8B, 0xFA)   # violeta/lavanda
C_CYAN     = RGBColor(0x00, 0xD2, 0xFF)   # cyan eléctrico
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY     = RGBColor(0xC8, 0xBE, 0xE6)
C_OK       = RGBColor(0x22, 0xC5, 0x5E)   # verde
C_WARNING  = RGBColor(0xFB, 0xBF, 0x24)   # amarillo
C_ERROR    = RGBColor(0xEF, 0x44, 0x44)   # rojo
C_SECTION  = RGBColor(0x5B, 0x21, 0xB6)   # violeta oscuro

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUTPUT_DIR = Path(__file__).resolve().parents[1] / "output"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb_hex(r: RGBColor) -> str:
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"


def _add_rect(slide, left, top, width, height, fill: RGBColor, alpha: int = 255):
    """Add a filled rectangle (no border)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.line.fill.background()  # no border
    fill_elem = shape.fill
    fill_elem.solid()
    fill_elem.fore_color.rgb = fill
    return shape


def _add_textbox(
    slide,
    left, top, width, height,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = C_WHITE,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
    wrap: bool = True,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _set_slide_bg(slide, color: RGBColor = C_BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_slide(prs: Presentation) -> object:
    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    return slide


def _header_bar(slide, title: str, subtitle: str = "") -> None:
    """Top accent bar with title."""
    _add_rect(slide,
              left=0, top=0,
              width=SLIDE_W, height=Inches(1.15),
              fill=C_SECTION)
    _add_textbox(slide,
                 left=Inches(0.35), top=Inches(0.12),
                 width=Inches(12.0), height=Inches(0.55),
                 text=title, font_size=26, bold=True,
                 color=C_ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_textbox(slide,
                     left=Inches(0.35), top=Inches(0.67),
                     width=Inches(12.0), height=Inches(0.35),
                     text=subtitle, font_size=13,
                     color=C_GRAY, align=PP_ALIGN.LEFT)


def _footer(slide, text: str = "ecoRTA · M4 QC técnico · 2026") -> None:
    _add_rect(slide,
              left=0, top=Inches(7.2),
              width=SLIDE_W, height=Inches(0.3),
              fill=C_SECTION)
    _add_textbox(slide,
                 left=Inches(0.35), top=Inches(7.21),
                 width=Inches(12.5), height=Inches(0.28),
                 text=text, font_size=9,
                 color=C_GRAY, align=PP_ALIGN.LEFT)


def _bullet_block(
    slide,
    left, top, width,
    items: list[tuple[str, str, RGBColor]],  # (icon, text, color)
    font_size: int = 14,
    line_height_in: float = 0.42,
) -> None:
    """Render a vertical list of (icon, text) pairs."""
    for i, (icon, text, color) in enumerate(items):
        y = top + Inches(i * line_height_in)
        _add_textbox(slide,
                     left=left, top=y,
                     width=Inches(0.45), height=Inches(line_height_in),
                     text=icon, font_size=font_size + 2,
                     color=color, align=PP_ALIGN.CENTER)
        _add_textbox(slide,
                     left=left + Inches(0.48), top=y,
                     width=width - Inches(0.48), height=Inches(line_height_in),
                     text=text, font_size=font_size,
                     color=C_WHITE)


def _check_card(
    slide,
    left, top, width,
    code: str,
    title: str,
    thresholds: str,
    rationale: str,
    color: RGBColor = C_ACCENT,
) -> None:
    """Render a compact QC-check card."""
    card_h = Inches(1.55)
    # card background
    _add_rect(slide, left, top, width, card_h, fill=RGBColor(0x1A, 0x10, 0x35))
    # left accent stripe
    _add_rect(slide, left, top, Inches(0.07), card_h, fill=color)
    # code badge
    _add_textbox(slide,
                 left=left + Inches(0.14), top=top + Inches(0.05),
                 width=width - Inches(0.2), height=Inches(0.28),
                 text=code, font_size=9, bold=True,
                 color=color, align=PP_ALIGN.LEFT)
    # title
    _add_textbox(slide,
                 left=left + Inches(0.14), top=top + Inches(0.3),
                 width=width - Inches(0.2), height=Inches(0.32),
                 text=title, font_size=13, bold=True,
                 color=C_WHITE, align=PP_ALIGN.LEFT)
    # thresholds
    _add_textbox(slide,
                 left=left + Inches(0.14), top=top + Inches(0.6),
                 width=width - Inches(0.2), height=Inches(0.36),
                 text=thresholds, font_size=10,
                 color=C_WARNING, align=PP_ALIGN.LEFT)
    # rationale
    _add_textbox(slide,
                 left=left + Inches(0.14), top=top + Inches(0.95),
                 width=width - Inches(0.2), height=Inches(0.5),
                 text=rationale, font_size=10, italic=True,
                 color=C_GRAY, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs: Presentation) -> None:
    slide = _add_slide(prs)

    # Big background glow rectangle (decorative)
    _add_rect(slide,
              left=Inches(7.8), top=Inches(1.5),
              width=Inches(5.2), height=Inches(5.0),
              fill=RGBColor(0x1A, 0x10, 0x35))

    # Top micro-label
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(1.6),
                 width=Inches(8.0), height=Inches(0.35),
                 text="ecoRTA · Módulo M4 — Type Curve Overlay",
                 font_size=13, color=C_CYAN, align=PP_ALIGN.LEFT)

    # Main title
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(2.05),
                 width=Inches(7.8), height=Inches(1.8),
                 text="QC Técnico\ndel Match Visual",
                 font_size=42, bold=True,
                 color=C_ACCENT, align=PP_ALIGN.LEFT)

    # Subtitle
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(3.85),
                 width=Inches(7.8), height=Inches(0.45),
                 text="Diagnóstico automático de no-unicidad y calidad de datos",
                 font_size=16, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Author / context
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(5.2),
                 width=Inches(7.8), height=Inches(1.1),
                 text=(
                     "Robert E. Padrón García\n"
                     "Maestría en Ingeniería de Yacimientos\n"
                     "Fundación Universidad de América · 2026"
                 ),
                 font_size=12, color=C_GRAY, align=PP_ALIGN.LEFT)

    # Right panel: decorative stat badges
    for i, (val, label) in enumerate([
        ("6", "checks independientes"),
        ("3", "niveles de severidad"),
        ("199", "tests en verde"),
    ]):
        y = Inches(1.9 + i * 1.4)
        _add_rect(slide,
                  left=Inches(8.2), top=y,
                  width=Inches(4.5), height=Inches(1.15),
                  fill=RGBColor(0x0E, 0x08, 0x24))
        _add_textbox(slide,
                     left=Inches(8.2), top=y + Inches(0.08),
                     width=Inches(4.5), height=Inches(0.6),
                     text=val, font_size=36, bold=True,
                     color=C_ACCENT, align=PP_ALIGN.CENTER)
        _add_textbox(slide,
                     left=Inches(8.2), top=y + Inches(0.65),
                     width=Inches(4.5), height=Inches(0.4),
                     text=label, font_size=12,
                     color=C_GRAY, align=PP_ALIGN.CENTER)

    _footer(slide)


def slide_motivacion(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "¿Por qué QC técnico en el match visual?",
                "El match log-log de Fetkovich es inherentemente no-único — múltiples curvas pueden ajustarse al mismo dataset")
    _footer(slide)

    # Left column: problem statement
    _add_textbox(slide,
                 left=Inches(0.4), top=Inches(1.3),
                 width=Inches(5.8), height=Inches(0.35),
                 text="El problema", font_size=15, bold=True, color=C_ACCENT)

    problems = [
        ("⚠️", "Pocos puntos → match visual ambiguo"),
        ("⚠️", "BHP muy variable → q/Δp no representativa"),
        ("⚠️", "Rango MBT estrecho → multiplicador X indeterminado"),
        ("⚠️", "Multiplicadores en 1.0 → match no realizado"),
        ("⚠️", "Todo el dataset en zona qDd ≈ 1 → Y no anclada"),
        ("⚠️", "Solo flujo transiente visible → kh y N no separables"),
    ]
    _bullet_block(slide,
                  left=Inches(0.4), top=Inches(1.75),
                  width=Inches(5.8),
                  items=[(ic, tx, C_WARNING) for ic, tx in problems],
                  font_size=13, line_height_in=0.44)

    # Right column: solution
    _add_textbox(slide,
                 left=Inches(6.9), top=Inches(1.3),
                 width=Inches(6.0), height=Inches(0.35),
                 text="La solución en ecoRTA", font_size=15, bold=True, color=C_CYAN)

    solutions = [
        ("✅", "6 checks automáticos en cada render del overlay"),
        ("✅", "Severidades ok / warning / error con umbrales de ingeniería"),
        ("✅", "Panel visual en la UI con detalle técnico accionable"),
        ("✅", "Basado en criterios Fetkovich (SPE-4629) y Palacio-Blasingame (SPE-25909)"),
        ("✅", "Diseñado para pozos exploratorios CPO-9 / Llanos Orientales"),
    ]
    _bullet_block(slide,
                  left=Inches(6.9), top=Inches(1.75),
                  width=Inches(6.0),
                  items=[(ic, tx, C_OK) for ic, tx in solutions],
                  font_size=13, line_height_in=0.44)

    # Ref quote
    _add_rect(slide,
              left=Inches(0.4), top=Inches(6.3),
              width=Inches(12.5), height=Inches(0.75),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(6.35),
                 width=Inches(12.1), height=Inches(0.6),
                 text=(
                     '"The type-curve matching process is not unique and can lead to erroneous results '
                     'if the diagnostics are ignored."  — Fetkovich, SPE-4629 (1980)'
                 ),
                 font_size=11, italic=True, color=C_GRAY)


def slide_arquitectura(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "Arquitectura del servicio",
                "rta_qc_service.py — separación estricta entre lógica de QC y UI")
    _footer(slide)

    # Flow: input → checks → aggregator → UI
    boxes = [
        (Inches(0.4),  Inches(2.0), Inches(2.4), Inches(1.1), "Inputs", "RTATransformPoint[ ]\nx_mult, y_mult", C_CYAN),
        (Inches(3.3),  Inches(1.3), Inches(5.6), Inches(2.5), "6 checks independientes",
         "POINT_COUNT · DRAWDOWN_STABILITY\nDATA_SPAN · MATCH_NOT_ADJUSTED\nQDD_RANGE · TRANSIENT_ONLY", C_ACCENT),
        (Inches(9.3),  Inches(1.8), Inches(2.4), Inches(1.4), "run_rta_qc()\naggregator", "list[QCResult]", C_ACCENT),
        (Inches(9.3),  Inches(3.5), Inches(2.4), Inches(1.1), "qc_severity_level()", "ok · warning · error", C_WARNING),
        (Inches(9.3),  Inches(5.0), Inches(3.2), Inches(0.9), "Expander QC técnico", "params_col en Streamlit", C_OK),
    ]

    for left, top, width, height, title, body, color in boxes:
        _add_rect(slide, left, top, width, height, fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, left, top, Inches(0.07), height, fill=color)
        _add_textbox(slide,
                     left=left + Inches(0.14), top=top + Inches(0.05),
                     width=width - Inches(0.2), height=Inches(0.32),
                     text=title, font_size=12, bold=True, color=color)
        _add_textbox(slide,
                     left=left + Inches(0.14), top=top + Inches(0.38),
                     width=width - Inches(0.2), height=height - Inches(0.4),
                     text=body, font_size=11, color=C_GRAY)

    # Arrows (text approximation)
    arrows = [
        (Inches(2.85), Inches(2.4), "→"),
        (Inches(8.95), Inches(2.4), "→"),
        (Inches(10.5), Inches(3.25), "↓"),
        (Inches(10.5), Inches(4.7), "↓"),
    ]
    for lft, tp, arrow in arrows:
        _add_textbox(slide,
                     left=lft, top=tp,
                     width=Inches(0.4), height=Inches(0.4),
                     text=arrow, font_size=20, bold=True,
                     color=C_ACCENT, align=PP_ALIGN.CENTER)

    # QCResult dataclass
    _add_rect(slide,
              left=Inches(0.4), top=Inches(5.5),
              width=Inches(8.5), height=Inches(1.4),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(5.55),
                 width=Inches(8.1), height=Inches(0.32),
                 text="@dataclass(frozen=True) class QCResult", font_size=11, bold=True, color=C_CYAN)
    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(5.88),
                 width=Inches(8.1), height=Inches(0.85),
                 text="    code: str          # identificador único del check\n"
                      "    severity: Literal['ok', 'warning', 'error']\n"
                      "    title: str          # título corto para la UI\n"
                      "    detail: str         # explicación técnica accionable",
                 font_size=10, color=C_GRAY)


def slide_checks_overview(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "Los 6 checks — vista general",
                "Cada check retorna exactamente un QCResult con la severidad correspondiente")
    _footer(slide)

    checks = [
        ("POINT_COUNT",         "Cantidad de puntos",          "error < 5 · warning < 15",        C_ERROR),
        ("DRAWDOWN_STABILITY",  "Estabilidad del drawdown",    "error CV > 30 % · warning > 15 %",C_WARNING),
        ("DATA_SPAN",           "Rango de MBT",                "error < 0.5 ciclos · warning < 1.0", C_WARNING),
        ("MATCH_NOT_ADJUSTED",  "Match realizado",             "warning si X = Y = 1.0",          C_WARNING),
        ("QDD_RANGE",           "Rango de qDd en overlay",     "warning si qDd_min > 0.70",       C_WARNING),
        ("TRANSIENT_ONLY",      "Flujo dominante",             "warning si pendiente log-log > −0.35", C_WARNING),
    ]

    cols = 3
    card_w = Inches(4.2)
    card_gap = Inches(0.17)
    for i, (code, title, thresh, color) in enumerate(checks):
        col = i % cols
        row = i // cols
        left = Inches(0.35) + col * (card_w + card_gap)
        top = Inches(1.3) + row * Inches(1.65)
        _check_card(slide, left, top, card_w, code, title, thresh,
                    _check_rationale_short(code), color)


def _check_rationale_short(code: str) -> str:
    return {
        "POINT_COUNT":        "< 5 puntos: el match no tiene base estadística.",
        "DRAWDOWN_STABILITY": "CV(Δp) alto → q/Δp dispersa → match ambiguo.",
        "DATA_SPAN":          "< 1 ciclo en MBT → multiplicador X no anclado.",
        "MATCH_NOT_ADJUSTED": "Multiplicadores en 1.0 → kh/N no tienen sentido.",
        "QDD_RANGE":          "Todos puntos en qDd ≈ 1 → Y sin información.",
        "TRANSIENT_ONLY":     "Sin BDF visible → kh y N no separables.",
    }.get(code, "")


def slide_check_point_count(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "POINT_COUNT — Cantidad de puntos válidos",
                "Primer filtro: si no hay datos suficientes, el resto del QC no tiene sentido")
    _footer(slide)

    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(1.35),
                 width=Inches(5.5), height=Inches(0.35),
                 text="Umbrales (configurables)", font_size=14, bold=True, color=C_ACCENT)

    thresholds = [
        (C_ERROR,   "🔴 error",   "n < 5 puntos",  "Sin base para el overlay — verificar filtros y columnas"),
        (C_WARNING, "⚠️  warning", "n < 15 puntos", "Match con alta incertidumbre — idealmente > 20 puntos"),
        (C_OK,      "✅ ok",       "n ≥ 15 puntos", "Suficientes para un overlay confiable"),
    ]
    for i, (color, sev, cond, msg) in enumerate(thresholds):
        y = Inches(1.85 + i * 1.1)
        _add_rect(slide, Inches(0.5), y, Inches(11.8), Inches(0.95),
                  fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, Inches(0.5), y, Inches(0.07), Inches(0.95), fill=color)
        _add_textbox(slide,
                     left=Inches(0.7), top=y + Inches(0.05),
                     width=Inches(1.5), height=Inches(0.35),
                     text=sev, font_size=12, bold=True, color=color)
        _add_textbox(slide,
                     left=Inches(2.3), top=y + Inches(0.05),
                     width=Inches(2.8), height=Inches(0.35),
                     text=cond, font_size=12, color=C_WHITE, bold=True)
        _add_textbox(slide,
                     left=Inches(0.7), top=y + Inches(0.48),
                     width=Inches(11.2), height=Inches(0.4),
                     text=msg, font_size=11, color=C_GRAY, italic=True)

    # Implementation note
    _add_rect(slide, Inches(0.5), Inches(5.4), Inches(11.8), Inches(1.4),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide,
                 left=Inches(0.7), top=Inches(5.45),
                 width=Inches(11.3), height=Inches(0.32),
                 text="Referencia técnica", font_size=12, bold=True, color=C_CYAN)
    _add_textbox(slide,
                 left=Inches(0.7), top=Inches(5.78),
                 width=Inches(11.3), height=Inches(0.9),
                 text=(
                     "Fetkovich (SPE-4629) recomienda al menos 15-20 puntos para anclar visualmente "
                     "la nube en las zonas transiente y BDF simultáneamente. Con < 5 puntos, la curva "
                     "tipo seleccionada (re/rw) y el multiplicador X son prácticamente arbitrarios."
                 ),
                 font_size=11, color=C_GRAY, italic=True)


def slide_check_drawdown(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "DRAWDOWN_STABILITY — Estabilidad del BHP de fondo",
                "La tasa normalizada q/Δp solo es un proxy de qDd si Δp es aproximadamente constante")
    _footer(slide)

    # Formula
    _add_rect(slide, Inches(0.5), Inches(1.3), Inches(11.8), Inches(1.05),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide,
                 left=Inches(0.7), top=Inches(1.35),
                 width=Inches(11.3), height=Inches(0.35),
                 text="Métrica utilizada: Coeficiente de Variación de Δp = CV(Δp) = σ(Δp) / μ(Δp)",
                 font_size=13, bold=True, color=C_CYAN)
    _add_textbox(slide,
                 left=Inches(0.7), top=Inches(1.72),
                 width=Inches(11.3), height=Inches(0.5),
                 text="donde  Δp = pi − pwf_used   [psia]   y   la media es el promedio aritmético de todos los puntos válidos",
                 font_size=11, color=C_GRAY, italic=True)

    thresholds = [
        (C_ERROR,   "🔴 error",   "CV > 30 %", "La tasa normalizada no es representativa — considerar estabilizar BHP o usar Δp promedio"),
        (C_WARNING, "⚠️  warning", "CV > 15 %", "Dispersión moderada — verificar que pwf_used sea el BHP representativo del período"),
        (C_OK,      "✅ ok",       "CV ≤ 15 %", "Drawdown estable — la normalización q/Δp es confiable"),
    ]
    for i, (color, sev, cond, msg) in enumerate(thresholds):
        y = Inches(2.5 + i * 1.0)
        _add_rect(slide, Inches(0.5), y, Inches(11.8), Inches(0.88),
                  fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, Inches(0.5), y, Inches(0.07), Inches(0.88), fill=color)
        _add_textbox(slide,
                     left=Inches(0.7), top=y + Inches(0.05),
                     width=Inches(1.5), height=Inches(0.32),
                     text=sev, font_size=12, bold=True, color=color)
        _add_textbox(slide,
                     left=Inches(2.3), top=y + Inches(0.05),
                     width=Inches(2.0), height=Inches(0.32),
                     text=cond, font_size=12, bold=True, color=C_WHITE)
        _add_textbox(slide,
                     left=Inches(0.7), top=y + Inches(0.43),
                     width=Inches(11.2), height=Inches(0.38),
                     text=msg, font_size=11, color=C_GRAY, italic=True)

    # Context note
    _add_rect(slide, Inches(0.5), Inches(5.65), Inches(11.8), Inches(1.25),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide,
                 left=Inches(0.7), top=Inches(5.7),
                 width=Inches(11.3), height=Inches(0.32),
                 text="Contexto CPO-9 / Llanos Orientales", font_size=12, bold=True, color=C_CYAN)
    _add_textbox(slide,
                 left=Inches(0.7), top=Inches(6.03),
                 width=Inches(11.3), height=Inches(0.78),
                 text=(
                     "En pruebas extensas con ESP, el BHP puede variar significativamente durante los ciclos de operación. "
                     "Un CV > 15 % es frecuente en crudos pesados (9-16 °API) donde la viscosidad alta "
                     "genera caídas de presión variables. Se recomienda usar el BHP estabilizado del período más largo."
                 ),
                 font_size=11, color=C_GRAY, italic=True)


def slide_check_span_and_match(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "DATA_SPAN · MATCH_NOT_ADJUSTED",
                "Rango de datos en el eje X y verificación de que el match fue realizado")
    _footer(slide)

    # DATA_SPAN card (left)
    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(1.3),
                 width=Inches(6.0), height=Inches(0.35),
                 text="DATA_SPAN — Ciclos logarítmicos de MBT",
                 font_size=14, bold=True, color=C_ACCENT)

    span_items = [
        (C_ERROR,   "🔴", "< 0.5 ciclos",  "Multiplicador X prácticamente libre — no aporta información sobre re/rw"),
        (C_WARNING, "⚠️",  "0.5 – 1.0 ciclos", "Match X con alta incertidumbre"),
        (C_OK,      "✅", "≥ 1.0 ciclos",  "Buen rango para anclar el multiplicador X"),
    ]
    for i, (color, icon, cond, msg) in enumerate(span_items):
        y = Inches(1.75 + i * 1.2)
        _add_rect(slide, Inches(0.5), y, Inches(6.2), Inches(1.0),
                  fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, Inches(0.5), y, Inches(0.07), Inches(1.0), fill=color)
        _add_textbox(slide, Inches(0.65), y + Inches(0.07),
                     Inches(1.3), Inches(0.32),
                     text=f"{icon} {cond}", font_size=11, bold=True, color=color)
        _add_textbox(slide, Inches(0.65), y + Inches(0.45),
                     Inches(5.9), Inches(0.48),
                     text=msg, font_size=10, color=C_GRAY, italic=True)

    # MATCH_NOT_ADJUSTED card (right)
    _add_textbox(slide,
                 left=Inches(7.0), top=Inches(1.3),
                 width=Inches(6.0), height=Inches(0.35),
                 text="MATCH_NOT_ADJUSTED — ¿Se realizó el match?",
                 font_size=14, bold=True, color=C_ACCENT)

    match_items = [
        (C_WARNING, "⚠️", "X = 1.0  Y = 1.0",
         "Ambos multiplicadores en posición inicial → kh, k y N no tienen significado físico"),
        (C_WARNING, "⚠️", "X ajustado · Y = 1.0",
         "kh no puede estimarse — ajuste la escala Y para anclar verticalmente la nube"),
        (C_OK,      "✅", "Y ajustado",
         "Match realizado — los parámetros son interpretables"),
    ]
    for i, (color, icon, cond, msg) in enumerate(match_items):
        y = Inches(1.75 + i * 1.4)
        _add_rect(slide, Inches(7.0), y, Inches(6.0), Inches(1.2),
                  fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, Inches(7.0), y, Inches(0.07), Inches(1.2), fill=color)
        _add_textbox(slide, Inches(7.15), y + Inches(0.07),
                     Inches(5.7), Inches(0.32),
                     text=f"{icon} {cond}", font_size=11, bold=True, color=color)
        _add_textbox(slide, Inches(7.15), y + Inches(0.5),
                     Inches(5.7), Inches(0.58),
                     text=msg, font_size=10, color=C_GRAY, italic=True)

    # Tolerance note
    _add_textbox(slide,
                 left=Inches(7.0), top=Inches(6.1),
                 width=Inches(6.0), height=Inches(0.55),
                 text="Tolerancia: |mult − 1.0| < 0.001 → se considera posición inicial",
                 font_size=10, color=C_WARNING, italic=True)


def slide_check_qdd_transient(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "QDD_RANGE · TRANSIENT_ONLY",
                "No-unicidad del eje Y y ausencia de BDF en los datos")
    _footer(slide)

    # QDD_RANGE (left)
    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(1.3),
                 width=Inches(6.2), height=Inches(0.35),
                 text="QDD_RANGE — Rango de qDd en el overlay",
                 font_size=14, bold=True, color=C_ACCENT)
    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(1.7),
                 width=Inches(6.2), height=Inches(0.7),
                 text=(
                     "qDd = normalized_rate × y_multiplier\n"
                     "Si todos los puntos tienen qDd > 0.70, el dataset solo cubre\n"
                     "la parte alta de la curva tipo → el multiplicador Y no está anclado."
                 ),
                 font_size=11, color=C_GRAY, italic=True)

    _add_rect(slide, Inches(0.5), Inches(2.55), Inches(6.2), Inches(0.95),
              fill=RGBColor(0x1A, 0x10, 0x35))
    _add_rect(slide, Inches(0.5), Inches(2.55), Inches(0.07), Inches(0.95), fill=C_WARNING)
    _add_textbox(slide, Inches(0.65), Inches(2.6),
                 Inches(5.9), Inches(0.32),
                 text="⚠️  warning  |  qDd_min > 0.70",
                 font_size=12, bold=True, color=C_WARNING)
    _add_textbox(slide, Inches(0.65), Inches(2.95),
                 Inches(5.9), Inches(0.45),
                 text="Obtener datos de mayor duración para anclar el match vertical.",
                 font_size=10, color=C_GRAY, italic=True)

    _add_rect(slide, Inches(0.5), Inches(3.6), Inches(6.2), Inches(0.95),
              fill=RGBColor(0x1A, 0x10, 0x35))
    _add_rect(slide, Inches(0.5), Inches(3.6), Inches(0.07), Inches(0.95), fill=C_OK)
    _add_textbox(slide, Inches(0.65), Inches(3.65),
                 Inches(5.9), Inches(0.32),
                 text="✅ ok  |  qDd_min ≤ 0.70",
                 font_size=12, bold=True, color=C_OK)
    _add_textbox(slide, Inches(0.65), Inches(4.0),
                 Inches(5.9), Inches(0.45),
                 text="Datos con suficiente declinación — match Y confiable.",
                 font_size=10, color=C_GRAY, italic=True)

    # TRANSIENT_ONLY (right)
    _add_textbox(slide,
                 left=Inches(7.0), top=Inches(1.3),
                 width=Inches(6.0), height=Inches(0.35),
                 text="TRANSIENT_ONLY — Flujo dominante en los últimos puntos",
                 font_size=14, bold=True, color=C_ACCENT)
    _add_textbox(slide,
                 left=Inches(7.0), top=Inches(1.7),
                 width=Inches(6.0), height=Inches(0.7),
                 text=(
                     "Algoritmo: regresión lineal en log-log del último 25 % de puntos\n"
                     "(mínimo 3). Pendiente m = d(log q/Δp) / d(log MBT).\n"
                     "BDF exponencial → m ≈ −1.0.  Transiente IARF → m ≈ −0.5."
                 ),
                 font_size=11, color=C_GRAY, italic=True)

    slopes = [
        (C_WARNING, "⚠️  warning", "m > −0.35", "Flujo transiente o lineal — BDF no alcanzado; kh y N no separables"),
        (C_OK,      "✅ ok",       "m ≤ −0.35", "Señal de BDF detectada — match interpretable"),
    ]
    for i, (color, sev, cond, msg) in enumerate(slopes):
        y = Inches(2.55 + i * 1.15)
        _add_rect(slide, Inches(7.0), y, Inches(6.1), Inches(0.95),
                  fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, Inches(7.0), y, Inches(0.07), Inches(0.95), fill=color)
        _add_textbox(slide, Inches(7.15), y + Inches(0.07),
                     Inches(5.8), Inches(0.32),
                     text=f"{sev}  ·  {cond}", font_size=12, bold=True, color=color)
        _add_textbox(slide, Inches(7.15), y + Inches(0.5),
                     Inches(5.8), Inches(0.38),
                     text=msg, font_size=10, color=C_GRAY, italic=True)

    # Palacio-Blasingame reference
    _add_rect(slide, Inches(7.0), Inches(4.95), Inches(6.1), Inches(0.9),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide, Inches(7.15), Inches(5.0),
                 Inches(5.8), Inches(0.8),
                 text=(
                     "Criterio de BDF: Palacio-Blasingame (SPE-25909) — cuando qDdi ≈ qDdid "
                     "en el gráfico log-log, el pozo está en pseudo-estado estacionario."
                 ),
                 font_size=10, color=C_GRAY, italic=True)

    # Limitar check note
    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(5.0),
                 width=Inches(6.2), height=Inches(0.55),
                 text="⚠  Check omitido si y_mult ≈ 1.0 (match no realizado) o n < 4 puntos",
                 font_size=10, color=C_WARNING, italic=True)


def slide_ui_integration(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "Integración en la UI Streamlit",
                "El QC se calcula en cada render del overlay — sin costo adicional para el usuario")
    _footer(slide)

    # Left: flow description
    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(1.3),
                 width=Inches(5.8), height=Inches(0.35),
                 text="Flujo de datos", font_size=14, bold=True, color=C_ACCENT)

    flow = [
        ("1", "CSV cargado con qo_stb_d + pwf_used_psia", C_CYAN),
        ("2", "compute_rta_transforms() → method_points", C_CYAN),
        ("3", "_qc_transform_points = method_points  [right_col]", C_CYAN),
        ("4", "run_rta_qc(points, x_mult, y_mult)  [params_col]", C_ACCENT),
        ("5", "qc_severity_level() → ✅ / ⚠️ / 🔴", C_ACCENT),
        ("6", "Expander auto-expandido si severity ≠ ok", C_OK),
    ]
    for i, (num, txt, color) in enumerate(flow):
        y = Inches(1.75 + i * 0.72)
        _add_rect(slide, Inches(0.5), y, Inches(0.42), Inches(0.55),
                  fill=RGBColor(0x1A, 0x10, 0x35))
        _add_textbox(slide, Inches(0.5), y + Inches(0.07),
                     Inches(0.42), Inches(0.42),
                     text=num, font_size=13, bold=True,
                     color=color, align=PP_ALIGN.CENTER)
        _add_textbox(slide, Inches(1.05), y + Inches(0.07),
                     Inches(5.2), Inches(0.42),
                     text=txt, font_size=11, color=C_WHITE)

    # Right: mock UI screenshot description
    _add_rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.5),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide, Inches(7.0), Inches(1.35),
                 Inches(5.8), Inches(0.35),
                 text="Panel params_col (lado derecho del overlay)",
                 font_size=11, bold=True, color=C_CYAN)

    # Mock UI components
    mock = [
        ("kh (mD·ft)",    "247.5",   C_WHITE),
        ("k (mD)",        "4.950",   C_WHITE),
        ("N vol. (MM STB)", "12.340", C_WHITE),
        ("─" * 40,        "",        C_SECTION),
        ("⚠️ QC técnico  [expandido]", "", C_WARNING),
        ("  🔴 Drawdown muy inestable", "", C_ERROR),
        ("  ⚠️ Rango de MBT limitado", "", C_WARNING),
        ("  ✅ Cantidad de puntos",     "", C_OK),
        ("─" * 40,         "",         C_SECTION),
        ("📌 Guardar match",  "",       C_ACCENT),
    ]
    for i, (label, val, color) in enumerate(mock):
        y = Inches(1.8 + i * 0.44)
        _add_textbox(slide, Inches(7.1), y,
                     Inches(3.8), Inches(0.4),
                     text=label, font_size=10, color=color)
        if val:
            _add_textbox(slide, Inches(10.9), y,
                         Inches(1.8), Inches(0.4),
                         text=val, font_size=10, bold=True,
                         color=C_WHITE, align=PP_ALIGN.RIGHT)

    # Footer note
    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(6.5),
                 width=Inches(6.0), height=Inches(0.5),
                 text="El expander se colapsa automáticamente cuando todos los checks son ✅ ok",
                 font_size=10, color=C_GRAY, italic=True)


def slide_tests(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "Cobertura de tests — 40 nuevos tests",
                "tests/test_rta_qc_service.py · pytest pasa en verde con 199 tests totales")
    _footer(slide)

    test_groups = [
        ("TestCheckPointCount",        "6 tests", "Boundaries: error en 3/4 pts, warning 5-14, ok ≥ 15"),
        ("TestCheckDrawdownStability", "5 tests", "CV=0 → ok; CV≈28 % → warning; CV alto → error; Δp=0 → error"),
        ("TestCheckDataSpan",          "5 tests", "2.0 ciclos → ok; 0.7 ciclos → warning; 0.2 ciclos → error"),
        ("TestCheckMatchAdjusted",     "6 tests", "Tolerancia 1e-3; X=Y=1 → warning; solo Y=1 → warning; Y ajust → ok"),
        ("TestCheckQddRange",          "5 tests", "y_mult=1.0 → skip; qDd_min>0.70 → warning; min_qdd en detail"),
        ("TestCheckTransientOnly",     "4 tests", "< 4 pts → []; slope≈-1 → ok; slope=0 → warning"),
        ("TestRunRtaQc",               "4 tests", "Aggregator: lista QCResult, orden correcto, vacío → POINT_COUNT"),
        ("TestQcSeverityLevel",        "5 tests", "Precedencia: error > warning > ok; lista vacía → ok"),
    ]

    col_gap = Inches(0.15)
    card_w = (SLIDE_W - Inches(0.5) * 2 - col_gap * 3) / 4 - Inches(0.01)

    for i, (cls, count, desc) in enumerate(test_groups):
        col = i % 4
        row = i // 4
        left = Inches(0.5) + col * (card_w + col_gap)
        top = Inches(1.35) + row * Inches(2.2)
        card_h = Inches(2.0)

        _add_rect(slide, left, top, card_w, card_h, fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, left, top, Inches(0.07), card_h, fill=C_OK)

        _add_textbox(slide,
                     left=left + Inches(0.14), top=top + Inches(0.07),
                     width=card_w - Inches(0.2), height=Inches(0.28),
                     text=count, font_size=22, bold=True, color=C_OK)
        _add_textbox(slide,
                     left=left + Inches(0.14), top=top + Inches(0.38),
                     width=card_w - Inches(0.2), height=Inches(0.38),
                     text=cls.replace("Test", ""), font_size=10, bold=True, color=C_WHITE)
        _add_textbox(slide,
                     left=left + Inches(0.14), top=top + Inches(0.82),
                     width=card_w - Inches(0.2), height=Inches(1.05),
                     text=desc, font_size=9, color=C_GRAY, italic=True)

    # Total badge
    _add_rect(slide, Inches(5.0), Inches(5.95), Inches(3.33), Inches(1.2),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide, Inches(5.0), Inches(6.0),
                 Inches(3.33), Inches(0.55),
                 text="199 tests", font_size=26, bold=True,
                 color=C_ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(5.0), Inches(6.55),
                 Inches(3.33), Inches(0.45),
                 text="en verde · 0 fallos · 0 skips",
                 font_size=11, color=C_OK, align=PP_ALIGN.CENTER)


def slide_non_uniqueness(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "Marco teórico: no-unicidad del match visual",
                "Fetkovich (1980) — por qué el QC es necesario, no opcional")
    _footer(slide)

    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(1.3),
                 width=Inches(12.3), height=Inches(0.4),
                 text="Fuentes de no-unicidad identificadas por Fetkovich (SPE-4629)",
                 font_size=15, bold=True, color=C_ACCENT)

    sources = [
        ("Eje X (MBT)",
         "El multiplicador X escala el tiempo adimensional tDd. Sin al menos 1 ciclo logarítmico "
         "de datos, múltiples combinaciones (re/rw, x_mult) dan el mismo ajuste visual. "
         "→ DATA_SPAN detecta esta condición.",
         C_WARNING),
        ("Eje Y (qDd)",
         "El multiplicador Y mapea q/Δp → qDd. Si todos los datos están en qDd > 0.70, "
         "la curva tipo podría ser cualquier b ∈ [0,1] con diferente kh. "
         "→ QDD_RANGE detecta esta condición.",
         C_WARNING),
        ("Zona transiente pura",
         "En flujo transiente (log-log slope ≈ 0 a -0.5), todas las curvas tipo se superponen. "
         "El match de kh y N simultáneamente requiere BDF visible. "
         "→ TRANSIENT_ONLY detecta esta condición.",
         C_WARNING),
        ("BHP variable",
         "Si Δp cambia significativamente, q/Δp no equivale a qDd. El match visual puede ser "
         "correcto pero la inferencia de kh será errónea. "
         "→ DRAWDOWN_STABILITY cuantifica este riesgo.",
         C_WARNING),
    ]

    for i, (title, desc, color) in enumerate(sources):
        y = Inches(1.85 + i * 1.18)
        _add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.05),
                  fill=RGBColor(0x1A, 0x10, 0x35))
        _add_rect(slide, Inches(0.5), y, Inches(0.07), Inches(1.05), fill=color)
        _add_textbox(slide, Inches(0.65), y + Inches(0.05),
                     Inches(2.5), Inches(0.32),
                     text=title, font_size=12, bold=True, color=C_WHITE)
        _add_textbox(slide, Inches(0.65), y + Inches(0.42),
                     Inches(11.8), Inches(0.55),
                     text=desc, font_size=10, color=C_GRAY, italic=True)


def slide_summary(prs: Presentation) -> None:
    slide = _add_slide(prs)
    _header_bar(slide,
                "Resumen y próximos pasos",
                "Estado actual del QC técnico M4 en ecoRTA")
    _footer(slide)

    # Left: what was done
    _add_textbox(slide,
                 left=Inches(0.5), top=Inches(1.3),
                 width=Inches(5.8), height=Inches(0.35),
                 text="Implementado en este commit", font_size=14, bold=True, color=C_OK)

    done = [
        "✅  rta_qc_service.py — 6 checks + aggregator + severity helper",
        "✅  test_rta_qc_service.py — 40 tests, 8 clases, todos en verde",
        "✅  UI integrada en m4_type_curve_overlay.py (expander params_col)",
        "✅  Auto-expansión cuando severity ≠ ok",
        "✅  Solo activo cuando CSV tiene columnas RTA válidas",
        "✅  Commit + push a feature/m4-type-curve-overlay",
    ]
    _bullet_block(slide,
                  left=Inches(0.5), top=Inches(1.78),
                  width=Inches(5.9),
                  items=[("", t, C_OK) for t in done],
                  font_size=12, line_height_in=0.52)

    # Right: next steps
    _add_textbox(slide,
                 left=Inches(7.0), top=Inches(1.3),
                 width=Inches(5.8), height=Inches(0.35),
                 text="Próximos commits pendientes", font_size=14, bold=True, color=C_ACCENT)

    pending = [
        "M1 — Módulo Pwf (BHP calculation)",
        "M2 — Correlaciones PVT (Standing / Beggs-Robinson)",
        "M5.1 — Modelo común de resultados M1-M4",
        "M5.2 — Dashboard comparativo EUR vs OOIP",
        "M5.3 — Reporte exportable (CSV/JSON/Excel)",
        "M5.4 — QC final y trazabilidad",
        "M5.5 — Tabla tesis vs Harmony (IHS/Fekete)",
    ]
    _bullet_block(slide,
                  left=Inches(7.0), top=Inches(1.78),
                  width=Inches(5.9),
                  items=[("→", t, C_ACCENT) for t in pending],
                  font_size=12, line_height_in=0.52)

    # Bottom: validation milestone
    _add_rect(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
              fill=RGBColor(0x0E, 0x08, 0x24))
    _add_textbox(slide, Inches(0.7), Inches(6.15),
                 Inches(11.8), Inches(0.75),
                 text=(
                     "Hito de validación: kh recuperado con error < 0.1 % usando caso sintético de declinación exponencial "
                     "(Fetkovich BDF, b=0, Pwf constante). Verificado por 18 tests en test_rta_synthetic_validation.py. "
                     "Los 6 checks de QC garantizan que el ingeniero detecte condiciones de no-unicidad antes de reportar resultados."
                 ),
                 font_size=10, color=C_GRAY, italic=True)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_DIR / "ecoRTA_QC_tecnico_M4.pptx"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_motivacion(prs)
    slide_non_uniqueness(prs)
    slide_arquitectura(prs)
    slide_checks_overview(prs)
    slide_check_point_count(prs)
    slide_check_drawdown(prs)
    slide_check_span_and_match(prs)
    slide_check_qdd_transient(prs)
    slide_ui_integration(prs)
    slide_tests(prs)
    slide_summary(prs)

    prs.save(str(out_path))
    print(f"[OK] Presentacion guardada en: {out_path}")
    print(f"     {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
